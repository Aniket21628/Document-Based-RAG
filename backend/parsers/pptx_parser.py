from pptx import Presentation
from typing import Dict
import logging

logger = logging.getLogger(__name__)

class PPTXParser:
    @staticmethod
    def parse(file_path: str) -> Dict[str, any]:
        """Parse PPTX and extract text content"""
        try:
            prs = Presentation(file_path)
            
            text_content = ""
            slide_count = 0
            
            for slide_num, slide in enumerate(prs.slides):
                slide_count += 1
                text_content += f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"
                    
                    # Extract table content
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            text_content += row_text + "\n"
            
            metadata = {
                "num_slides": slide_count,
                "file_type": "pptx",
                "file_path": file_path
            }
            
            return {
                "content": text_content,
                "metadata": metadata
            }
            
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {str(e)}")
            raise Exception(f"Failed to parse PPTX: {str(e)}")