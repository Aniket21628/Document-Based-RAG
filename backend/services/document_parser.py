import os
from typing import List, Dict, Any
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
import logging

logger = logging.getLogger(__name__)

class DocumentParser:
    
    @staticmethod
    def parse_pdf(file_path: str) -> List[str]:
        """Parse PDF and return text chunks"""
        chunks = []
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                for page_num, page in enumerate(pdf_reader.pages):
                    text = page.extract_text()
                    if text.strip():
                        chunks.append(f"Page {page_num + 1}: {text}")
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {e}")
            raise
        return chunks
    
    @staticmethod
    def parse_docx(file_path: str) -> List[str]:
        """Parse DOCX and return text chunks"""
        chunks = []
        try:
            doc = Document(file_path)
            current_chunk = ""
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    current_chunk += paragraph.text + "\n"
                    if len(current_chunk) > 1000:  # Chunk size
                        chunks.append(current_chunk.strip())
                        current_chunk = ""
            
            if current_chunk.strip():
                chunks.append(current_chunk.strip())
                
        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {e}")
            raise
        return chunks
    
    @staticmethod
    def parse_pptx(file_path: str) -> List[str]:
        """Parse PPTX and return text chunks"""
        chunks = []
        try:
            prs = Presentation(file_path)
            for slide_num, slide in enumerate(prs.slides):
                slide_text = f"Slide {slide_num + 1}:\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text:
                        slide_text += shape.text + "\n"
                if slide_text.strip() != f"Slide {slide_num + 1}:":
                    chunks.append(slide_text.strip())
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {e}")
            raise
        return chunks
    
    @staticmethod
    def parse_csv(file_path: str) -> List[str]:
        """Parse CSV and return text chunks"""
        chunks = []
        try:
            df = pd.read_csv(file_path)
            
            # Create header information chunk
            header_info = f"CSV Headers: {', '.join(df.columns.tolist())}\n"
            header_info += f"Total Rows: {len(df)}\n"
            chunks.append(header_info)
            
            # Create chunks from rows
            chunk_size = 50  # rows per chunk
            for i in range(0, len(df), chunk_size):
                chunk_df = df.iloc[i:i+chunk_size]
                chunk_text = f"Rows {i+1}-{min(i+chunk_size, len(df))}:\n"
                chunk_text += chunk_df.to_string(index=False)
                chunks.append(chunk_text)
                
        except Exception as e:
            logger.error(f"Error parsing CSV {file_path}: {e}")
            raise
        return chunks
    
    @staticmethod
    def parse_txt(file_path: str) -> List[str]:
        """Parse TXT/MD files and return text chunks"""
        chunks = []
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
                
            # Split into chunks of ~1000 characters
            chunk_size = 1000
            for i in range(0, len(content), chunk_size):
                chunk = content[i:i+chunk_size]
                if chunk.strip():
                    chunks.append(chunk.strip())
                    
        except Exception as e:
            logger.error(f"Error parsing TXT {file_path}: {e}")
            raise
        return chunks
    
    @classmethod
    def parse_document(cls, file_path: str, file_type: str) -> List[str]:
        """Main method to parse any supported document type"""
        parsers = {
            'pdf': cls.parse_pdf,
            'docx': cls.parse_docx,
            'pptx': cls.parse_pptx,
            'csv': cls.parse_csv,
            'txt': cls.parse_txt,
            'md': cls.parse_txt
        }
        
        if file_type.lower() not in parsers:
            raise ValueError(f"Unsupported file type: {file_type}")
        
        return parsers[file_type.lower()](file_path)