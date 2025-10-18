import os
import logging
from typing import Dict, Any
import PyPDF2
from docx import Document
from pptx import Presentation
import pandas as pd
import markdown
import re

logger = logging.getLogger(__name__)

class DocumentProcessor:
    """Process different document types"""
    
    def process_file(self, file_path: str) -> Dict[str, Any]:
        """Process a file based on its extension"""
        ext = os.path.splitext(file_path)[1].lower()
        
        processors = {
            '.pdf': self._process_pdf,
            '.docx': self._process_docx,
            '.pptx': self._process_pptx,
            '.csv': self._process_csv,
            '.txt': self._process_txt,
            '.md': self._process_md,
        }
        
        processor = processors.get(ext)
        if not processor:
            raise ValueError(f"Unsupported file type: {ext}")
        
        return processor(file_path)
    
    def _process_pdf(self, file_path: str) -> Dict[str, Any]:
        """Process PDF file"""
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                text_content = ""
                
                for page_num, page in enumerate(pdf_reader.pages):
                    page_text = page.extract_text()
                    text_content += f"\n--- Page {page_num + 1} ---\n{page_text}"
                
                return {
                    "content": text_content,
                    "metadata": {
                        "num_pages": len(pdf_reader.pages),
                        "file_type": "pdf",
                        "file_path": file_path
                    }
                }
        except Exception as e:
            logger.error(f"Error parsing PDF {file_path}: {str(e)}")
            raise
    
    def _process_docx(self, file_path: str) -> Dict[str, Any]:
        """Process DOCX file"""
        try:
            doc = Document(file_path)
            text_content = ""
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content += paragraph.text + "\n"
            
            for table in doc.tables:
                for row in table.rows:
                    row_text = " | ".join([cell.text for cell in row.cells])
                    text_content += row_text + "\n"
            
            return {
                "content": text_content,
                "metadata": {
                    "num_paragraphs": len(doc.paragraphs),
                    "num_tables": len(doc.tables),
                    "file_type": "docx",
                    "file_path": file_path
                }
            }
        except Exception as e:
            logger.error(f"Error parsing DOCX {file_path}: {str(e)}")
            raise
    
    def _process_pptx(self, file_path: str) -> Dict[str, Any]:
        """Process PPTX file"""
        try:
            prs = Presentation(file_path)
            text_content = ""
            
            for slide_num, slide in enumerate(prs.slides):
                text_content += f"\n--- Slide {slide_num + 1} ---\n"
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"
                    
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text for cell in row.cells])
                            text_content += row_text + "\n"
            
            return {
                "content": text_content,
                "metadata": {
                    "num_slides": len(prs.slides),
                    "file_type": "pptx",
                    "file_path": file_path
                }
            }
        except Exception as e:
            logger.error(f"Error parsing PPTX {file_path}: {str(e)}")
            raise
    
    def _process_csv(self, file_path: str) -> Dict[str, Any]:
        """Process CSV file"""
        try:
            df = pd.read_csv(file_path)
            
            text_content = f"CSV Data Summary:\n"
            text_content += f"Columns: {', '.join(df.columns.tolist())}\n"
            text_content += f"Number of rows: {len(df)}\n\n"
            
            text_content += "Column Information:\n"
            for col in df.columns:
                text_content += f"- {col}: {df[col].dtype}\n"
            
            text_content += "\nData Preview:\n"
            text_content += df.head(10).to_string(index=False)
            
            numeric_cols = df.select_dtypes(include=['number']).columns
            if len(numeric_cols) > 0:
                text_content += "\n\nNumerical Statistics:\n"
                text_content += df[numeric_cols].describe().to_string()
            
            return {
                "content": text_content,
                "metadata": {
                    "num_rows": len(df),
                    "num_columns": len(df.columns),
                    "columns": df.columns.tolist(),
                    "file_type": "csv",
                    "file_path": file_path
                }
            }
        except Exception as e:
            logger.error(f"Error parsing CSV {file_path}: {str(e)}")
            raise
    
    def _process_txt(self, file_path: str) -> Dict[str, Any]:
        """Process TXT file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            return {
                "content": content,
                "metadata": {
                    "file_type": "txt",
                    "file_path": file_path,
                    "character_count": len(content),
                    "word_count": len(content.split())
                }
            }
        except Exception as e:
            logger.error(f"Error parsing TXT {file_path}: {str(e)}")
            raise
    
    def _process_md(self, file_path: str) -> Dict[str, Any]:
        """Process Markdown file"""
        try:
            with open(file_path, 'r', encoding='utf-8') as file:
                content = file.read()
            
            # Convert markdown to plain text
            html = markdown.markdown(content)
            plain_text = re.sub('<[^<]+?>', '', html)
            
            return {
                "content": plain_text,
                "metadata": {
                    "file_type": "markdown",
                    "file_path": file_path,
                    "character_count": len(plain_text),
                    "word_count": len(plain_text.split())
                }
            }
        except Exception as e:
            logger.error(f"Error parsing Markdown {file_path}: {str(e)}")
            raise